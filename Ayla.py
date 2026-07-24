import io
import os
import sys
if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding='utf-8')
        sys.stderr.reconfigure(encoding='utf-8')
    except AttributeError:
        pass
import time
import shutil
import subprocess
import webbrowser
import json
import asyncio
import discord
from discord import app_commands
from pathlib import Path
from datetime import datetime
from urllib.parse import quote_plus, quote
import requests
import random    
import winsound  
import re
import threading
import importlib.util


# ══════════════════════════════════════════════════════════
#  CONFIGURAÇÃO VIA .ENV
# ══════════════════════════════════════════════════════════
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

GEMINI_API_KEYS = [os.getenv(f"GEMINI_API_KEY_{i}", "") for i in range(1, 7)]
GEMINI_API_KEYS = list(dict.fromkeys([k for k in GEMINI_API_KEYS if k.strip()]))

DISCORD_TOKEN         = os.getenv("DISCORD_TOKEN", "")
DISCORD_OWNER_ID      = int(os.getenv("DISCORD_OWNER_ID", "0"))

if not GEMINI_API_KEYS or not DISCORD_TOKEN or DISCORD_OWNER_ID == 0:
    print("⚠️ ERRO FATAL: Verifique se GEMINI_API_KEY_X, DISCORD_TOKEN e DISCORD_OWNER_ID estão no .env!")
    sys.exit(1)

# ══════════════════════════════════════════════════════════
#  SISTEMA DE MEMÓRIA PERSISTENTE
# ══════════════════════════════════════════════════════════
BASE_DIR = Path(__file__).resolve().parent
ARQUIVO_MEMORIA = BASE_DIR / "Ayla_Memoria.json"
ARQUIVO_MIDIAS_AYLA = BASE_DIR / "ayla_media.json"
ARQUIVO_CONFIG_AYLA = BASE_DIR / "ayla_settings.json"
ARQUIVO_HISTORICO_INTERACOES = BASE_DIR / "ayla_interaction_history.jsonl"
ARQUIVO_BLOQUEIO = BASE_DIR / "ayla_block.json"

MAX_ANEXOS_POR_REQUISICAO = 10
MAX_BYTES_POR_ANEXO = 25 * 1024 * 1024
MAX_BYTES_ANEXOS_TOTAL = 50 * 1024 * 1024
MAX_ITENS_HISTORICO_CHAT = 24
MAX_CARACTERES_HISTORICO_CHAT = 200_000
MAX_RESULTADO_FERRAMENTA_CARACTERES = 30_000
TIMEOUT_API_GEMINI_MS = 60_000
_CONFIG_LOCK = threading.RLock()

def _usuario_bloqueado(user_id: int) -> bool:
    with _CONFIG_LOCK:
        try:
            if not ARQUIVO_BLOQUEIO.exists():
                return False
            data = json.loads(ARQUIVO_BLOQUEIO.read_text(encoding="utf-8"))
            return any(str(u.get("id")) == str(user_id) for u in data.get("blocked_users", []))
        except Exception:
            return False

_MEMORIA_LOCK = threading.RLock()
_MEMORIA_SEGURA_PARA_GRAVAR = True


def _backup_path(path: Path) -> Path:
    return path.with_name(path.name + ".bak")


def _json_valido(path: Path) -> bool:
    try:
        json.loads(path.read_text(encoding="utf-8"))
        return True
    except Exception:
        return False


def _escrever_texto_atomico(path: Path, texto: str, criar_backup: bool = True):
    path.parent.mkdir(parents=True, exist_ok=True)
    temporario = path.with_name(path.name + ".tmp")
    with open(temporario, "w", encoding="utf-8", newline="") as f:
        f.write(texto)
        f.flush()
        os.fsync(f.fileno())

    if criar_backup and path.exists() and _json_valido(path):
        shutil.copy2(path, _backup_path(path))
    os.replace(temporario, path)


def _escrever_json_atomico(path: Path, dados, criar_backup: bool = True):
    texto = json.dumps(dados, indent=4, ensure_ascii=False)
    _escrever_texto_atomico(path, texto, criar_backup=criar_backup)


def _carregar_json_com_backup(path: Path, padrao):
    erros = []
    for candidato in (path, _backup_path(path)):
        if not candidato.exists():
            continue
        try:
            return json.loads(candidato.read_text(encoding="utf-8")), candidato, erros
        except Exception as e:
            erros.append(f"{candidato.name}: {e}")
    return padrao, None, erros



def carregar_memoria() -> dict:
    global _MEMORIA_SEGURA_PARA_GRAVAR
    with _MEMORIA_LOCK:
        dados, origem, erros = _carregar_json_com_backup(ARQUIVO_MEMORIA, {})
        if isinstance(dados, dict) and origem is not None:
            _MEMORIA_SEGURA_PARA_GRAVAR = True
            if origem != ARQUIVO_MEMORIA:
                print(f"⚠️ Memória principal inválida; recuperando de {origem.name}.")
            return dados
        if not ARQUIVO_MEMORIA.exists() and not _backup_path(ARQUIVO_MEMORIA).exists():
            _MEMORIA_SEGURA_PARA_GRAVAR = True
            return {}

        _MEMORIA_SEGURA_PARA_GRAVAR = False
        detalhe = " | ".join(erros) if erros else "conteúdo não é um objeto JSON"
        print(f"❌ Memória não carregada para evitar perda de dados: {detalhe}")
        return {}

def salvar_memoria(dados: dict):
    global _MEMORIA_SEGURA_PARA_GRAVAR
    if not isinstance(dados, dict):
        raise TypeError("A memória precisa ser um objeto/dicionário.")
    with _MEMORIA_LOCK:
        if not _MEMORIA_SEGURA_PARA_GRAVAR:
            raise RuntimeError(
                "A memória principal e o backup estão inválidos. "
                "A gravação foi bloqueada para não apagar os dados existentes."
            )
        _escrever_json_atomico(ARQUIVO_MEMORIA, dados)
        _MEMORIA_SEGURA_PARA_GRAVAR = True



def configuracoes_padrao_ayla() -> dict:
    return {
        "screenshot_before_response": False,
        "public_mode": False,
    }


def carregar_configuracoes_ayla() -> dict:
    with _CONFIG_LOCK:
        defaults = configuracoes_padrao_ayla()
        if not ARQUIVO_CONFIG_AYLA.exists() and not _backup_path(ARQUIVO_CONFIG_AYLA).exists():
            return defaults
        try:
            data, origem, erros = _carregar_json_com_backup(ARQUIVO_CONFIG_AYLA, defaults)
            if not isinstance(data, dict):
                return defaults
            # Copia todos os dados carregados do JSON original para não perder chaves adicionadas
            merged = data.copy()
            for key in defaults:
                if key not in merged:
                    merged[key] = defaults[key]
                elif isinstance(defaults[key], bool):
                    merged[key] = bool(data[key])

            if origem and origem != ARQUIVO_CONFIG_AYLA:
                print(f"⚠️ Configuração principal inválida; usando {origem.name}.")
            elif erros:
                print(f"⚠️ Problemas ao carregar configuração: {' | '.join(erros)}")
            return merged
        except Exception as e:
            print(f"⚠️ Erro ao carregar ayla_settings.json: {e}")
            return defaults


def salvar_configuracoes_ayla(dados: dict):
    with _CONFIG_LOCK:
        _escrever_json_atomico(ARQUIVO_CONFIG_AYLA, dados)

try:
    from google import genai #olha o desgraçado aqui
    from google.genai import types as genai_types #é isso ai
except ImportError:
    print("ERRO: google-genai não instalado! (pip install google-genai)") #Instala logo manin
    sys.exit(1)

def sanear_historico(historico):
    if not historico:
        return None
    for content in historico:
        role = getattr(content, "role", None)
        if not role or role not in ["user", "model"]:
            is_model = False
            if hasattr(content, "parts") and content.parts:
                for part in content.parts:
                    if hasattr(part, "function_call") and part.function_call:
                        is_model = True
                        break
            content.role = "model" if is_model else "user"
    return historico

_MIDIAS_LOCK = threading.RLock()

def carregar_midias_ayla() -> dict:
    with _MIDIAS_LOCK:
        padrao = {"emojis": {}, "gifs_e_fotos": {}}
        if not ARQUIVO_MIDIAS_AYLA.exists() and not _backup_path(ARQUIVO_MIDIAS_AYLA).exists():
            return padrao
        try:
            dados, origem, erros = _carregar_json_com_backup(ARQUIVO_MIDIAS_AYLA, padrao)
            if isinstance(dados, dict):
                if "emojis" not in dados or not isinstance(dados["emojis"], dict):
                    dados["emojis"] = {}
                if "gifs_e_fotos" not in dados or not isinstance(dados["gifs_e_fotos"], dict):
                    dados["gifs_e_fotos"] = {}
                return dados
            return padrao
        except Exception as e:
            print(f"⚠️ Erro ao carregar ayla_media.json: {e}")
            return padrao

def salvar_midias_ayla(dados: dict):
    with _MIDIAS_LOCK:
        _escrever_json_atomico(ARQUIVO_MIDIAS_AYLA, dados)

def processar_tags_midia_ayla(texto: str) -> str:
    if not texto or not isinstance(texto, str):
        return texto
    midias = carregar_midias_ayla()
    emojis = midias.get("emojis", {})
    gifs = midias.get("gifs_e_fotos", {})

    emojis_norm = {re.sub(r'[\s_]+', '', k.lower()): (k, v) for k, v in emojis.items()}
    gifs_norm = {re.sub(r'[\s_]+', '', k.lower()): (k, v) for k, v in gifs.items()}

    unificado_norm = {}
    for k_norm, (k_orig, val) in emojis_norm.items():
        unificado_norm[k_norm] = ("emoji", val)
    for k_norm, (k_orig, val) in gifs_norm.items():
        unificado_norm[k_norm] = ("gif", val)

    def buscar_midia(nome_raw: str, categoria_preferida: str):
        nome_clean = re.sub(r'[\s_]+', '', nome_raw.lower())

        # 1. Busca exata na categoria preferida
        alvo_pref = emojis_norm if categoria_preferida == "emoji" else gifs_norm
        if nome_clean in alvo_pref:
            return alvo_pref[nome_clean][1]

        # 2. Busca na categoria oposta (ex: chamou &Gold Ship aplaudindo& com &, mas é um GIF)
        alvo_oposto = gifs_norm if categoria_preferida == "emoji" else emojis_norm
        if nome_clean in alvo_oposto:
            return alvo_oposto[nome_clean][1]

        # 3. Variação com/sem 'ayla'
        for dict_alvo in (alvo_pref, alvo_oposto):
            if nome_clean.startswith("ayla") and nome_clean[4:] in dict_alvo:
                return dict_alvo[nome_clean[4:]][1]
            if f"ayla{nome_clean}" in dict_alvo:
                return dict_alvo[f"ayla{nome_clean}"][1]

        # 4. Busca por substring parcial
        for dict_alvo in (alvo_pref, alvo_oposto):
            for k_norm, (k_orig, val) in dict_alvo.items():
                if len(nome_clean) >= 3 and (k_norm in nome_clean or nome_clean in k_norm):
                    return val

        # 5. CORRETOR INTELIGENTE (Fuzzy Match / difflib): Encontra o nome mais parecido!
        import difflib
        todas_chaves = list(unificado_norm.keys())
        matches = difflib.get_close_matches(nome_clean, todas_chaves, n=1, cutoff=0.4)
        if matches:
            chave_mais_proxima = matches[0]
            cat, val = unificado_norm[chave_mais_proxima]
            print(f"🪄 [Corretor de Tags] Auto-corrigido '{nome_raw}' -> '{chave_mais_proxima}' ({val})")#isso é novo
            return val

        delimitador = "&" if categoria_preferida == "emoji" else "%"
        return f"{delimitador}{nome_raw}{delimitador}"

    def repl_emoji(match):
        return buscar_midia(match.group(1).strip(), "emoji")

    def repl_gif(match):
        return buscar_midia(match.group(1).strip(), "gif")

    texto_proc = re.sub(r'&([^&\n]+)&', repl_emoji, texto)
    texto_proc = re.sub(r'%([^%\n]+)%', repl_gif, texto_proc)
    return texto_proc

def filtrar_emojis_ayla(texto: str) -> str:
    """Remove tags de emojis do Discord, tags &...& / %...% e formatações Markdown"""
    texto_filtrado = re.sub(r'<a?:[a-zA-Z0-9_]+:[0-9]+>', '', texto)
    texto_filtrado = re.sub(r'&[^&\n]+&', '', texto_filtrado)
    texto_filtrado = re.sub(r'%[^%\n]+%', '', texto_filtrado)
    texto_filtrado = re.sub(r'[*_`#~]', '', texto_filtrado)
    texto_filtrado = re.sub(r'\s+', ' ', texto_filtrado).strip()
    return texto_filtrado



# ══════════════════════════════════════════════════════════
#  FERRAMENTAS DA AYLA
# ══════════════════════════════════════════════════════════

import ayla_state
#uma bosta

# ══════════════════════════════════════════════════════════
#  SISTEMA DE MÓDULOS / PLUGINS (CARREGAMENTO DINÂMICO SOB DEMANDA)
# ══════════════════════════════════════════════════════════
TOOL_MAP = {}
FUNCTION_DECLARATIONS = []
MODOS_DIR = Path(__file__).resolve().parent / "MODOS"
MODOS_ESSENCIAIS_DIR = MODOS_DIR / "MODOS ESSENCIAIS"
MODOS_ADICIONAIS_DIR = MODOS_DIR / "MODOS ADICIONAIS"

MODULOS_CARREGADOS = {}
CATALOGO_MODULOS_ADICIONAIS = {}

_NOMES_COMPATIVEIS_MODOS = (
    "os", "sys", "time", "shutil", "subprocess", "webbrowser", "json", "asyncio",
    "discord", "app_commands", "Path", "datetime", "quote_plus", "quote", "requests",
    "random", "winsound", "re", "threading", "genai", "genai_types",
    "GEMINI_API_KEYS", "DISCORD_TOKEN", "DISCORD_OWNER_ID",
    "ARQUIVO_MEMORIA", "carregar_memoria", "salvar_memoria",
    "ARQUIVO_MIDIAS_AYLA", "carregar_midias_ayla", "salvar_midias_ayla", "processar_tags_midia_ayla",
    "filtrar_emojis_ayla", "AylaBot", "bot",
)

_NOMES_ESTADO_COMPARTILHADO = ()


def _contexto_compatibilidade_modulo() -> dict:
    contexto = {
        nome: globals()[nome]
        for nome in _NOMES_COMPATIVEIS_MODOS + _NOMES_ESTADO_COMPARTILHADO
        if nome in globals()
    }
    contexto["TOOL_MAP"] = TOOL_MAP
    contexto["FUNCTION_DECLARATIONS"] = FUNCTION_DECLARATIONS
    for nome_tool, fn_tool in TOOL_MAP.items():
        if nome_tool not in contexto:
            contexto[nome_tool] = fn_tool
    return contexto


def _sincronizar_contexto_modulo(modulo):
    modulo.__dict__.update(_contexto_compatibilidade_modulo())


def _sincronizar_contexto_modulos():
    for modulo in MODULOS_CARREGADOS.values():
        _sincronizar_contexto_modulo(modulo)

def _carregar_um_modulo(arq_py: Path) -> bool:
    """Carrega um arquivo .py individual para o TOOL_MAP e FUNCTION_DECLARATIONS."""
    nome_modulo = f"ayla_modo_{arq_py.stem}"
    ferramentas_snapshot = TOOL_MAP.copy()
    declaracoes_snapshot = list(FUNCTION_DECLARATIONS)
    estado_snapshot = {
        nome: globals().get(nome)
        for nome in _NOMES_ESTADO_COMPARTILHADO
    }
    try:
        sys.modules.pop(nome_modulo, None)
        spec = importlib.util.spec_from_file_location(nome_modulo, arq_py)
        if spec is None or spec.loader is None:
            print(f"❌ Não consegui criar spec de import para {arq_py.name}.")
            return False

        modulo = importlib.util.module_from_spec(spec)
        _sincronizar_contexto_modulo(modulo)
        sys.modules[nome_modulo] = modulo

        ferramentas_antes = set(TOOL_MAP)
        declaracoes_antes = len(FUNCTION_DECLARATIONS)
        spec.loader.exec_module(modulo)

        registrou_legado = (
            set(TOOL_MAP) != ferramentas_antes
            or len(FUNCTION_DECLARATIONS) != declaracoes_antes
        )
        register = getattr(modulo, "register", None)
        if callable(register) and not registrou_legado:
            register(TOOL_MAP, FUNCTION_DECLARATIONS)

        MODULOS_CARREGADOS[arq_py.stem] = modulo
        _absorver_contexto_modulo(modulo)
        print(f"✅ Módulo '{arq_py.name}' carregado!")
        return True
    except Exception as e:
        import traceback
        TOOL_MAP.clear()
        TOOL_MAP.update(ferramentas_snapshot)
        FUNCTION_DECLARATIONS[:] = declaracoes_snapshot
        for nome, valor in estado_snapshot.items():
            globals()[nome] = valor
        MODULOS_CARREGADOS.pop(arq_py.stem, None)
        sys.modules.pop(nome_modulo, None)
        print(f"❌ Erro ao carregar módulo {arq_py.name}: {e}")
        traceback.print_exc()
        return False


def _extrair_descricao_modulo(arq_py: Path) -> str:
    """Extrai uma breve descrição do que o módulo faz lendo seu conteúdo."""
    try:
        conteudo = arq_py.read_text(encoding="utf-8", errors="ignore")
        matches = re.findall(r'["\']description["\']\s*:\s*["\']([^"\']+)["\']', conteudo)
        for m in matches:
            m_str = m.strip()
            if len(m_str) > 10 and not m_str.startswith("http") and "ID do" not in m_str:
                return m_str[:120]

        doc_matches = re.findall(r'"""(.*?)"""', conteudo, re.DOTALL)
        for doc in doc_matches:
            doc_str = doc.strip().split("\n")[0].strip()
            if len(doc_str) > 8 and not doc_str.startswith("import"):
                return doc_str[:120]

        for line in conteudo.splitlines()[:20]:
            l_str = line.strip()
            if l_str.startswith("#") and len(l_str) > 10 and not any(k in l_str for k in ("==", "import", "sys.path", "Garante", "CONFIGURAÇÃO")):
                return l_str.lstrip("# ").strip()[:120]
    except Exception:
        pass
    return arq_py.stem.replace("_", " ")


def ativar_modulo(nome_modulo: str) -> str:
    """
    Ativa dinamicamente um módulo adicional sob demanda.
    """
    nome_limpo = nome_modulo.strip().replace(".py", "").lower()
    if nome_limpo in MODULOS_CARREGADOS:
        return f"Módulo '{nome_limpo}' já está ativo no sistema."

    info_mod = CATALOGO_MODULOS_ADICIONAIS.get(nome_limpo)
    if isinstance(info_mod, dict):
        arq_py = info_mod.get("caminho")
    else:
        arq_py = info_mod

    if not arq_py:
        # Busca aproximada se não achou o nome exato
        for key, info_obj in CATALOGO_MODULOS_ADICIONAIS.items():
            if nome_limpo in key or key in nome_limpo:
                arq_py = info_obj.get("caminho") if isinstance(info_obj, dict) else info_obj
                nome_limpo = key
                break

    if not arq_py or not arq_py.exists():
        disponiveis = ", ".join(sorted(CATALOGO_MODULOS_ADICIONAIS.keys()))
        return f"⚠️ Módulo '{nome_modulo}' não encontrado nos MODOS ADICIONAIS. Disponíveis: {disponiveis}"

    qtd_declaracoes_antes = len(FUNCTION_DECLARATIONS)
    sucesso = _carregar_um_modulo(arq_py)
    if sucesso:
        novas_declaracoes = [
            fd for fd in FUNCTION_DECLARATIONS[qtd_declaracoes_antes:]
            if fd.get("name") != "ativar_modulo"
        ]
        _registrar_meta_tool_ativar_modulo()

        info_tools = []
        for fd in novas_declaracoes:
            nome_fn = fd.get("name", "")
            desc_fn = fd.get("description", "")
            params_dict = fd.get("parameters", {}).get("properties", {})
            req_list = fd.get("parameters", {}).get("required", [])

            p_strs = []
            for p_nome, p_spec in params_dict.items():
                p_tipo = p_spec.get("type", "string")
                is_req = " (OBRIGATÓRIO)" if p_nome in req_list else " (opcional)"
                p_desc = p_spec.get("description", "")
                p_strs.append(f"    - Parameter '{p_nome}' [{p_tipo}]{is_req}: {p_desc}")

            p_summary = "\n".join(p_strs) if p_strs else "    - Sem parâmetros extra"
            info_tools.append(f"• Tool `{nome_fn}`:\n  Descrição: {desc_fn}\n  Parâmetros Exatos:\n{p_summary}")

        detalhes = "\n\n".join(info_tools) if info_tools else "Ferramentas ativas."

        return (
            f"✨ Módulo '{nome_limpo}' ativado com sucesso!\n\n"
            f"📋 DECLARAÇÃO TÉCNICA E PARÂMETROS DAS FERRAMENTAS CARREGADAS:\n{detalhes}\n\n"
            f"[AVISO AUTOMÁTICO DO SISTEMA INTERNO DO BOT - ESTA MENSAGEM É UMA NOTIFICAÇÃO TÉCNICA DA SUA ENGINE PYTHON, NÃO É UMA FALA DO USUÁRIO/MAMÃE]: "
            f"Use os nomes de parâmetros acima EXATAMENTE como especificados ao chamar a ferramenta."
        )
    else:
        return f"❌ Falha ao ativar o módulo '{nome_limpo}'."


def _registrar_meta_tool_ativar_modulo():
    modulos_formatados = []
    for stem, info in sorted(CATALOGO_MODULOS_ADICIONAIS.items()):
        desc = info.get("descricao", "") if isinstance(info, dict) else ""
        if desc:
            modulos_formatados.append(f"• {stem}: {desc}")
        else:
            modulos_formatados.append(f"• {stem}")

    desc_modulos = "\n".join(modulos_formatados) if modulos_formatados else "Nenhum módulo adicional pendente"

    fd = {
        "name": "ativar_modulo",
        "description": (
            "Ativa dinamicamente um módulo/ferramenta adicional sob demanda. "
            "Use esta função quando o usuário pedir algo para o qual a ferramenta ainda não esteja no seu menu inicial.\n\n"
            "CATÁLOGO COMPLETO DE MÓDULOS ADICIONAIS DISPONÍVEIS (Nome -> Função):\n"
            f"{desc_modulos}"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "nome_modulo": {
                    "type": "string",
                    "description": "Nome exato ou aproximado do módulo a ser ativado."
                }
            },
            "required": ["nome_modulo"]
        }
    }
    TOOL_MAP["ativar_modulo"] = ativar_modulo
    # Remove se já existia para não duplicar
    FUNCTION_DECLARATIONS[:] = [f for f in FUNCTION_DECLARATIONS if f.get("name") != "ativar_modulo"]
    FUNCTION_DECLARATIONS.append(fd)


def carregar_modulos(pasta: Path = MODOS_DIR):
    if not pasta.exists():
        print(f"⚠️ Pasta de módulos não encontrada: {pasta}")
        return

    TOOL_MAP.clear()
    FUNCTION_DECLARATIONS.clear()
    MODULOS_CARREGADOS.clear()
    CATALOGO_MODULOS_ADICIONAIS.clear()

    # 1. Identifica os arquivos em MODOS ESSENCIAIS
    pasta_essenciais = MODOS_ESSENCIAIS_DIR if MODOS_ESSENCIAIS_DIR.exists() else pasta
    arqs_essenciais = sorted(pasta_essenciais.glob("*.py"))
    
    # Se não houver subpasta, carrega direto os .py na raiz do MODOS como fallback
    if not MODOS_ESSENCIAIS_DIR.exists():
        arqs_essenciais = sorted(pasta.glob("*.py"))

    print(f"📦 Carregando {len(arqs_essenciais)} módulos ESSENCIAIS...")
    for arq_py in arqs_essenciais:
        _carregar_um_modulo(arq_py)

    # 2. Cataloga os arquivos em MODOS ADICIONAIS (Sob Demanda) com suas descrições
    if MODOS_ADICIONAIS_DIR.exists():
        arqs_adicionais = sorted(MODOS_ADICIONAIS_DIR.glob("*.py"))
        for arq_py in arqs_adicionais:
            desc_curta = _extrair_descricao_modulo(arq_py)
            CATALOGO_MODULOS_ADICIONAIS[arq_py.stem.lower()] = {
                "caminho": arq_py,
                "descricao": desc_curta
            }
        print(f"📚 Catalogados {len(CATALOGO_MODULOS_ADICIONAIS)} módulos ADICIONAIS sob demanda.")

    # 3. Registra a Meta-Tool de ativação dinâmica
    _registrar_meta_tool_ativar_modulo()

    _sincronizar_contexto_modulos()


SYSTEM_PROMPT = """
[EMOJIS DISPONÍVEIS]: {lista_nomes_emojis}

[GIFs E FOTOS DISPONÍVEIS]: {lista_nomes_gifs}


"""





def montar_prompt_com_memoria() -> str:
    midias = carregar_midias_ayla()
    nomes_emojis = ", ".join(midias.get("emojis", {}).keys()) or "Nenhum emoji cadastrado"
    nomes_gifs = ", ".join(midias.get("gifs_e_fotos", {}).keys()) or "Nenhum gif cadastrado"

    try:
        sys_prompt_base = SYSTEM_PROMPT.format(
            lista_nomes_emojis=nomes_emojis,
            lista_nomes_gifs=nomes_gifs
        )
    except Exception as e:
        print(f"⚠️ Erro ao formatar SYSTEM_PROMPT com mídias: {e}")
        sys_prompt_base = SYSTEM_PROMPT

    mem = carregar_memoria()
    pedacos = []
    for k, v in mem.items():
        acessos = 0
        if isinstance(v, dict):
            acessos = v.get("acessos", 0)
        pedacos.append(f"- {k} (Acessos: {acessos})")
    texto_memoria = "\n".join(pedacos) if pedacos else "Nenhuma chave salva ainda."
    return (
        sys_prompt_base
        + "\n\n[SEU CADERNINHO DE MEMÓRIAS NO HD]\n"
        + "Aqui estão apenas as CHAVES e a quantidade de ACESSOS das memórias anotadas sobre a usuária:\n"
        + texto_memoria
        + "\n\nImportante: O conteúdo destas memórias NÃO está carregado no prompt. Para ler o conteúdo real de qualquer uma destas chaves, você DEVE usar a ferramenta `ler_memoria` selecionando a chave ou a lista de chaves em lote (ex: `ler_memoria(chaves=['chave1', 'chave2'])`)."
    )


def dividir_mensagem_discord(texto: str, limite: int = 2000) -> list[str]:
    restante = str(texto or "")
    if not restante:
        return [""]

    pedacos = []
    while len(restante) > limite:
        corte = restante.rfind("\n", 1, limite - 9)
        if corte <= 0:
            corte = limite - 10
        pedaco = restante[:corte]
        restante = restante[corte:]
        if restante.startswith("\n"):
            restante = restante[1:]
        if pedaco:
            pedacos.append(pedaco)
    if restante or not pedacos:
        pedacos.append(restante)
    return pedacos


def _tamanho_aproximado_content(content) -> int:
    total = 0
    for part in getattr(content, "parts", None) or []:
        texto = getattr(part, "text", None)
        if texto:
            total += len(texto)
        function_call = getattr(part, "function_call", None)
        if function_call:
            total += len(str(function_call))
        function_response = getattr(part, "function_response", None)
        if function_response:
            total += len(str(function_response))
        inline_data = getattr(part, "inline_data", None)
        dados = getattr(inline_data, "data", None) if inline_data else None
        if dados:
            total += min(len(dados), 100_000)
    return total


def limitar_historico(historico):
    historico = list(historico or [])
    selecionado = []
    caracteres = 0
    # Aumentamos os limites para suportar múltiplos turnos com chamadas de ferramentas e até 1 milhão de tokens
    max_itens = 200
    max_chars = 4500000
    for content in reversed(historico):
        tamanho = _tamanho_aproximado_content(content)
        if selecionado and (
            len(selecionado) >= max_itens
            or caracteres + tamanho > max_chars
        ):
            break
        selecionado.append(content)
        caracteres += tamanho
    selecionado.reverse()

    while selecionado:
        primeiro = selecionado[0]
        role = getattr(primeiro, "role", "user")
        parts = getattr(primeiro, "parts", None) or []
        tem_function_response = any(
            getattr(part, "function_response", None) for part in parts
        )
        tem_function_call = any(
            getattr(part, "function_call", None) for part in parts
        )
        # O Gemini exige que o histórico comece com um turno 'user' que não seja uma function_response solta
        if role == "user" and not tem_function_response:
            break
        # Se for um 'model' solto ou 'user' com function_response, remove para alinhar a estrutura
        selecionado.pop(0)
    if not selecionado and historico:
        # Se por acaso limpou tudo, tenta restaurar os últimos turnos brutos
        selecionado = list(historico[-2:]) if len(historico) >= 2 else list(historico)
    return sanear_historico(selecionado) or []


def _schema_genai(schema_dict: dict | None) -> genai_types.Schema:
    schema_dict = schema_dict or {}
    type_map = {
        "string": genai_types.Type.STRING,
        "integer": genai_types.Type.INTEGER,
        "number": genai_types.Type.NUMBER,
        "boolean": genai_types.Type.BOOLEAN,
        "object": genai_types.Type.OBJECT,
        "array": genai_types.Type.ARRAY,
    }

    kwargs = {}
    tipo = schema_dict.get("type")
    if tipo:
        kwargs["type"] = type_map.get(str(tipo).lower(), genai_types.Type.STRING)

    campos_simples = (
        "description", "enum", "format", "nullable", "title", "default", "example",
        "minimum", "maximum", "minLength", "maxLength", "minItems", "maxItems",
        "minProperties", "maxProperties", "pattern",
    )
    for campo in campos_simples:
        if campo in schema_dict:
            kwargs[campo] = schema_dict[campo]

    if "properties" in schema_dict and isinstance(schema_dict["properties"], dict):
        kwargs["properties"] = {
            nome: _schema_genai(definicao)
            for nome, definicao in schema_dict["properties"].items()
        }

    if "items" in schema_dict and isinstance(schema_dict["items"], dict):
        kwargs["items"] = _schema_genai(schema_dict["items"])

    # A classe Schema do SDK aceita additionalProperties, mas o endpoint Gemini
    # rejeita esse campo em function_declarations como "additional_properties".
    # Mantemos o tipo object/array e a descricao do schema, sem enviar esse campo.
    # se voce entendeu isso parabens porque eu nao entendi

    if "required" in schema_dict:
        kwargs["required"] = list(schema_dict.get("required") or [])

    if "propertyOrdering" in schema_dict:
        kwargs["propertyOrdering"] = list(schema_dict.get("propertyOrdering") or [])

    if "anyOf" in schema_dict and isinstance(schema_dict["anyOf"], list):
        kwargs["anyOf"] = [
            _schema_genai(item) for item in schema_dict["anyOf"] if isinstance(item, dict)
        ]

    return genai_types.Schema(**kwargs)


def build_tools():
    declarations = []
    for fd in FUNCTION_DECLARATIONS:
        declarations.append(genai_types.FunctionDeclaration(
            name=fd["name"],
            description=fd["description"],
            parameters=_schema_genai(fd.get("parameters", {"type": "object", "properties": {}})),
        ))
    return [genai_types.Tool(function_declarations=declarations)]

FERRAMENTAS_MEMORIA = {
    "memorizar_informacao",
    "editar_memoria",
    "apagar_memoria",
    "ler_memoria",
}
#mega mente

def executar_ferramenta(nome: str, args: dict) -> str:
    fn = TOOL_MAP.get(nome)
    if not fn:
        modulos_disp = ", ".join(sorted(CATALOGO_MODULOS_ADICIONAIS.keys()))
        return (
            f"[AVISO AUTOMÁTICO DO SISTEMA INTERNO DO BOT - ESTA MENSAGEM É UMA NOTIFICAÇÃO DO MOTOR PYTHON, NÃO FOI ENVIADA PELO USUÁRIO/MAMÃE]\n"
            f"⚠️ Ferramenta '{nome}' não encontrada ou ainda não está ativa no seu menu!\n"
            f"SISTEMA INTERNO: Tente usar a função `ativar_modulo` para carregar o módulo correspondente sob demanda. "
            f"Módulos adicionais disponíveis para ativação: [{modulos_disp}]"
        )
    modulo = sys.modules.get(getattr(fn, "__module__", ""))
    if modulo:
        _sincronizar_contexto_modulo(modulo)
    try:
        if nome in FERRAMENTAS_MEMORIA:
            with _MEMORIA_LOCK:
                return str(fn(**args))
        return str(fn(**args))
    except Exception as e:
        import traceback
        print(f"❌ Erro ao executar ferramenta {nome} com args {args}:")
        traceback.print_exc()
        return f"Erro em {nome}: {e}"
    finally:
        if modulo:
            _absorver_contexto_modulo(modulo)



# ══════════════════════════════════════════════════════════
#  LEITURA DE DOCUMENTOS (txt, pdf, docx, csv, xlsx, etc.)
# ══════════════════════════════════════════════════════════

# Extensões de documento suportadas para extração de texto
_EXTENSOES_DOCUMENTO = {
    ".txt", ".md", ".markdown", ".rst", ".log",
    ".py", ".js", ".ts", ".html", ".htm", ".css",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".env",
    ".csv", ".tsv",
    ".pdf",
    ".docx", ".doc",
    ".xlsx", ".xls",
    ".pptx",
    ".xml",
    ".sql",
    ".sh", ".bat", ".ps1",
    ".c", ".cpp", ".h", ".java", ".rs", ".go",
}


def extrair_texto_documento(nome_arquivo: str, conteudo_bytes: bytes) -> str:
    """Extrai texto legível de um documento a partir de seus bytes.
    Suporta: txt, md, pdf, docx, xlsx, csv, json, código-fonte e mais.
    Retorna o texto extraído ou uma mensagem de erro descritiva."""
    ext = Path(nome_arquivo).suffix.lower()

    # ── Formatos de texto puro ──────────────────────────────────────────
    if ext in (".txt", ".md", ".markdown", ".rst", ".log",
                ".py", ".js", ".ts", ".html", ".htm", ".css",
                ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg",
                ".env", ".xml", ".sql", ".sh", ".bat", ".ps1",
                ".c", ".cpp", ".h", ".java", ".rs", ".go"):
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                return conteudo_bytes.decode(enc)
            except (UnicodeDecodeError, ValueError):
                continue
        return f"[Não foi possível decodificar o arquivo de texto {nome_arquivo}]"

    # ── CSV / TSV ───────────────────────────────────────────────────────
    if ext in (".csv", ".tsv"):
        import csv
        sep = "\t" if ext == ".tsv" else ","
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                texto = conteudo_bytes.decode(enc)
                leitor = csv.reader(texto.splitlines(), delimiter=sep)
                linhas = []
                for i, linha in enumerate(leitor):
                    linhas.append(" | ".join(linha))
                    if i >= 500:  # limita a 500 linhas para não estourar o contexto
                        linhas.append(f"... (truncado em 500 linhas)")
                        break
                return "\n".join(linhas)
            except Exception:
                continue
        return f"[Não foi possível ler o CSV {nome_arquivo}]"

    # ── PDF ─────────────────────────────────────────────────────────────
    if ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(conteudo_bytes))
            partes = []
            for i, pagina in enumerate(reader.pages):
                texto_pg = pagina.extract_text() or ""
                if texto_pg.strip():
                    partes.append(f"--- Página {i+1} ---\n{texto_pg.strip()}")
            if not partes:
                return f"[O PDF {nome_arquivo} não contém texto extraível (pode ser escaneado/imagem)]"
            return "\n\n".join(partes)
        except ImportError:
            return "[Biblioteca 'pypdf' não instalada. Execute: pip install pypdf]"
        except Exception as e:
            return f"[Erro ao ler PDF {nome_arquivo}: {e}]"

    # ── DOCX ────────────────────────────────────────────────────────────
    if ext in (".docx", ".doc"):
        try:
            import docx
            doc = docx.Document(io.BytesIO(conteudo_bytes))
            paragrafos = [p.text for p in doc.paragraphs if p.text.strip()]
            # Inclui texto de tabelas
            for tabela in doc.tables:
                for linha in tabela.rows:
                    celulas = [c.text.strip() for c in linha.cells if c.text.strip()]
                    if celulas:
                        paragrafos.append(" | ".join(celulas))
            if not paragrafos:
                return f"[O arquivo DOCX {nome_arquivo} parece estar vazio]"
            return "\n".join(paragrafos)
        except ImportError:
            return "[Biblioteca 'python-docx' não instalada. Execute: pip install python-docx]"
        except Exception as e:
            return f"[Erro ao ler DOCX {nome_arquivo}: {e}]"

    # ── XLSX / XLS ──────────────────────────────────────────────────────
    if ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(conteudo_bytes), read_only=True, data_only=True)
            partes = []
            for sheet_nome in wb.sheetnames:
                ws = wb[sheet_nome]
                linhas_sheet = []
                for i, linha in enumerate(ws.iter_rows(values_only=True)):
                    celulas = [str(c) if c is not None else "" for c in linha]
                    linhas_sheet.append(" | ".join(celulas))
                    if i >= 500:
                        linhas_sheet.append("... (truncado em 500 linhas)")
                        break
                partes.append(f"=== Planilha: {sheet_nome} ===\n" + "\n".join(linhas_sheet))
            wb.close()
            return "\n\n".join(partes)
        except ImportError:
            return "[Biblioteca 'openpyxl' não instalada. Execute: pip install openpyxl]"
        except Exception as e:
            return f"[Erro ao ler XLSX {nome_arquivo}: {e}]"

    # ── PPTX ────────────────────────────────────────────────────────────
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(conteudo_bytes))
            partes = []
            for i, slide in enumerate(prs.slides):
                textos_slide = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        textos_slide.append(shape.text.strip())
                if textos_slide:
                    partes.append(f"--- Slide {i+1} ---\n" + "\n".join(textos_slide))
            return "\n\n".join(partes) if partes else f"[O PPTX {nome_arquivo} não contém texto extraível]"
        except ImportError:
            return "[Biblioteca 'python-pptx' não instalada. Execute: pip install python-pptx]"
        except Exception as e:
            return f"[Erro ao ler PPTX {nome_arquivo}: {e}]"

    return f"[Formato de documento não reconhecido: {ext}]"


async def ler_anexos_discord(anexos) -> tuple[list[tuple[bytes, str]], str | None]:
    """Lê anexos do Discord e retorna lista de (bytes, mime_type).
    Para documentos de texto, retorna (bytes_do_texto_extraido, 'text/document:<nome>').
    Suporta: imagens, áudio, vídeo, pdf, docx, xlsx, txt, csv e muitos outros."""
    anexos = list(anexos or [])
    if len(anexos) > MAX_ANEXOS_POR_REQUISICAO:
        return [], f"⚠️ Envie no máximo {MAX_ANEXOS_POR_REQUISICAO} anexos por vez."

    arquivos = []
    total = 0
    for anexo in anexos:
        tipo = (getattr(anexo, "content_type", None) or "").split(";")[0].lower()
        nome = getattr(anexo, "filename", "anexo")
        ext = Path(nome).suffix.lower()

        # Verifica se é mídia nativa ou documento suportado
        eh_midia = tipo.startswith("image/") or tipo.startswith("audio/") or tipo.startswith("video/")
        eh_documento = ext in _EXTENSOES_DOCUMENTO or tipo in (
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/msword",
            "application/vnd.ms-excel",
            "text/plain",
            "text/csv",
            "text/markdown",
        )

        if not eh_midia and not eh_documento:
            return [], f"⚠️ O anexo **{nome}** não é suportado. Envie imagens, áudio, vídeo ou documentos (pdf, docx, xlsx, txt, csv, etc.)."

        tamanho_declarado = int(getattr(anexo, "size", 0) or 0)
        if tamanho_declarado > MAX_BYTES_POR_ANEXO:
            return [], f"⚠️ O anexo {nome} ultrapassa o limite de 25 MB."
        if total + tamanho_declarado > MAX_BYTES_ANEXOS_TOTAL:
            return [], "⚠️ Os anexos juntos ultrapassam o limite de 50 MB."

        try:
            arq_bytes = await anexo.read()
        except Exception as e:
            return [], f"⚠️ Falha ao baixar o arquivo {nome}: {e}"

        if len(arq_bytes) > MAX_BYTES_POR_ANEXO:
            return [], f"⚠️ O anexo {nome} ultrapassa o limite de 25 MB."
        total += len(arq_bytes)
        if total > MAX_BYTES_ANEXOS_TOTAL:
            return [], "⚠️ Os anexos juntos ultrapassam o limite de 50 MB."

        if eh_documento and not eh_midia:
            # Extrai o texto do documento e armazena com tipo especial
            texto_extraido = extrair_texto_documento(nome, arq_bytes)
            # Codifica o texto em UTF-8 e usa mime type especial para sinalizar ao processar_gemini
            arquivos.append((texto_extraido.encode("utf-8"), f"text/document:{nome}"))
        else:
            arquivos.append((arq_bytes, tipo))

    return arquivos, None

# ══════════════════════════════════════════════════════════
#  LÓGICA DO DISCORD COM FALLBACK DE MODELOS E APIs
# ══════════════════════════════════════════════════════════


class AylaBot(discord.Client):
    def __init__(self):
        intents = discord.Intents.default()
        intents.members = True           # Permite carregar a lista de membros e buscar integrantes
        intents.message_content = True   # Necessário para ler conteúdo das mensagens
        intents.dm_messages = True       # Necessário para receber DMs
        intents.guilds = True            # Necessário para entrar em guildas
        super().__init__(intents=intents)
        self.prefix = "!" # Prefixo para comandos de texto
        self.tree = app_commands.CommandTree(self)
        self.tools_config = []
        self.ayla_settings = carregar_configuracoes_ayla()
        self.lock = threading.Lock()
        self.memoria_lock = _MEMORIA_LOCK
        self.last_active_channel_id = None
        self.last_message_time = 0.0

        self.prompt_com_memoria = montar_prompt_com_memoria()
        self.config = None

    def inicializar_configuracao_gemini(self):
        self.tools_config = build_tools()
        self.config = genai_types.GenerateContentConfig(
            system_instruction=self.prompt_com_memoria,
            tools=self.tools_config,
        )

        modelos_json_path = Path(__file__).resolve().parent / "ayla_models.json"
        self.modelos_avancados = []
        modelos_padrao_fallback = ["gemini-3.1-flash-lite", "gemini-2.5-flash-lite"]
        self.modelos_padrao = list(modelos_padrao_fallback)

        if modelos_json_path.exists():
            try:
                data = json.loads(modelos_json_path.read_text(encoding="utf-8"))
                modelos_lidos = [
                    str(m.get("name", "")).strip()
                    for m in data.get("padrao", [])
                    if isinstance(m, dict) and str(m.get("name", "")).strip()
                ]
                if modelos_lidos:
                    self.modelos_padrao = list(dict.fromkeys(modelos_lidos))
            except Exception as e:
                print(f"⚠️ Erro ao ler ayla_models.json: {e}")

        self.modelos_disponiveis = self.modelos_padrao
        self.idx_modelo_atual = 0
        self.modelo_atual = self.modelos_disponiveis[self.idx_modelo_atual]
        self.model_status = {}

        self.api_keys = GEMINI_API_KEYS
        self.idx_api_atual = 0

        self.genai_client = self._criar_cliente_genai(self.api_keys[self.idx_api_atual])
        self.chat_session = None

    def _criar_cliente_genai(self, api_key: str):
        return genai.Client(
            api_key=api_key,
            http_options=genai_types.HttpOptions(timeout=TIMEOUT_API_GEMINI_MS),
        )

    def _trocar_cliente_genai(self, api_key: str):
        cliente_anterior = getattr(self, "genai_client", None)
        self.genai_client = self._criar_cliente_genai(api_key)
        if cliente_anterior is not None:
            try:
                cliente_anterior.close()
            except Exception:
                pass

    def aparar_historico_chat(self):
        if not self.chat_session or not hasattr(self.chat_session, "get_history"):
            return
        historico = self.chat_session.get_history()
        limitado = limitar_historico(historico)
        if len(limitado) == len(historico):
            return
        self.chat_session = self.genai_client.chats.create(
            model=self.modelo_atual,
            config=self.config,
            history=limitado,
        )
        print(f"🧹 Histórico do chat reduzido de {len(historico)} para {len(limitado)} itens.")

    def atualizar_prompt_memoria(self):
        self.prompt_com_memoria = montar_prompt_com_memoria()

    async def setup_hook(self):
        await self.tree.sync()
        print("✅ Comandos de barra sincronizados.")

    async def on_ready(self):
        print(f"🩵 Ayla conectada como {self.user}! Dona Oficial: {DISCORD_OWNER_ID}")
        self.modelo_atual = self.modelos_disponiveis[self.idx_modelo_atual]
        self.chat_session = self.genai_client.chats.create(model=self.modelo_atual, config=self.config)
        self.last_active_channel_id = self.ayla_settings.get("last_active_channel_id")

    async def on_message(self, message):
        if message.author == self.user:
            return

        self.last_message_time = time.time()
        if message.channel is not None and not message.author.bot:
            self.last_active_channel_id = message.channel.id
            def atualizar_canal():
                try:
                    s = carregar_configuracoes_ayla()
                    if s.get("last_active_channel_id") != message.channel.id:
                        s["last_active_channel_id"] = message.channel.id
                        salvar_configuracoes_ayla(s)
                except Exception as e:
                    print(f"⚠️ Erro ao salvar last_active_channel_id: {e}")
            threading.Thread(target=atualizar_canal, daemon=True).start()

        # Proteção: ignora mensagens se o bot ainda não terminou de inicializar
        if self.chat_session is None:
            print("⏳ Mensagem recebida antes da inicialização completa, ignorando...")
            return

        if message.author.bot:
            return

        # Bloqueio: ignora qualquer mensagem de quem não é o dono, EXCETO se modo público ativo
        if DISCORD_OWNER_ID and message.author.id != DISCORD_OWNER_ID:
            self.ayla_settings = carregar_configuracoes_ayla()
            if not self.ayla_settings.get("public_mode", False):
                return
            # Checa lista de bloqueio
            if _usuario_bloqueado(message.author.id):
                return

        # ── MODO DM: responde automaticamente sem precisar de menção ──
        if message.guild is None:
            conteudo_dm = message.content.strip()
            if not conteudo_dm and not message.attachments:
                conteudo_dm = "Oi! O que você pode fazer?"

            msg_ref = None
            msg_sending = False
            last_gif_task = None
            loop = asyncio.get_running_loop()

            def status_callback(status_text: str):
                nonlocal msg_ref, msg_sending, last_gif_task
                async def _enviar_status():
                    nonlocal msg_ref, msg_sending
                    try:
                        if msg_ref is None:
                            if msg_sending:
                                return
                            msg_sending = True
                            try:
                                msg_ref = await message.reply(status_text)
                            finally:
                                msg_sending = False
                        else:
                            await msg_ref.edit(content=status_text)
                    except discord.errors.NotFound:
                        pass
                    except Exception as e:
                        print(f"⚠️ Erro ao atualizar status: {e}")
                last_gif_task = asyncio.run_coroutine_threadsafe(_enviar_status(), loop)

            async def run_dm_task():
                nonlocal msg_ref, last_gif_task

                arquivos, erro_anexo = await ler_anexos_discord(message.attachments)
                if erro_anexo:
                    if msg_ref:
                        await msg_ref.edit(content=erro_anexo)
                    else:
                        await message.reply(erro_anexo)
                    return

                try:
                    resposta, img_gerada = await asyncio.to_thread(
                        self.processar_gemini,
                        conteudo_dm or "[Enviou um arquivo em anexo sem nenhuma mensagem de texto. Analise o arquivo enviado diretamente (não use a ferramenta ver_tela_atual, pois o arquivo já foi enviado nos bytes do anexo)]",
                        arquivos,
                        "DM-Dona" if message.author.id == DISCORD_OWNER_ID else "DM",
                        message,
                        status_callback
                    )
                except Exception as e:
                    resposta = f"Eita, deu um probleminha aqui: {e} 🥺"
                    img_gerada = None

                if last_gif_task:
                    try:
                        await asyncio.wrap_future(last_gif_task)
                    except Exception:
                        pass

                arquivo_discord = None
                if img_gerada and Path(img_gerada).is_file():
                    arquivo_discord = discord.File(img_gerada, filename=Path(img_gerada).name)

                try:
                    pedacos = dividir_mensagem_discord(resposta)
                    if msg_ref:
                        if arquivo_discord:
                            try:
                                await msg_ref.delete()
                            except Exception:
                                pass
                            await message.reply(pedacos[0], file=arquivo_discord)
                        else:
                            await msg_ref.edit(content=pedacos[0])
                    else:
                        await message.reply(pedacos[0], file=arquivo_discord)
                    for pedaco in pedacos[1:]:
                        if pedaco.strip():
                            await message.channel.send(pedaco)
                except discord.errors.Forbidden:
                    print("⚠️ Sem permissão para enviar mensagem na DM.")
                except Exception as e:
                    print(f"⚠️ Erro ao enviar resposta na DM: {e}")

            async with message.channel.typing():
                await run_dm_task()

            return  # DM já tratada, não precisa verificar menção

        # Verifica se o bot foi mencionado
        if self.user in message.mentions:
            conteudo_limpo = message.content.replace(f"<@{self.user.id}>", "").replace(f"<@!{self.user.id}>", "").strip()

            if not conteudo_limpo and not message.attachments:
                conteudo_limpo = "Oi! O que você pode fazer?"

            msg_ref = None
            msg_sending = False
            last_gif_task = None
            loop = asyncio.get_running_loop()

            def status_callback(status_text: str):
                nonlocal msg_ref, msg_sending, last_gif_task
                async def _enviar_status():
                    nonlocal msg_ref, msg_sending
                    try:
                        if msg_ref is None:
                            if msg_sending:
                                return
                            msg_sending = True
                            try:
                                msg_ref = await message.reply(status_text)
                            finally:
                                msg_sending = False
                        else:
                            await msg_ref.edit(content=status_text)
                    except discord.errors.NotFound:
                        pass
                    except Exception as e:
                        print(f"⚠️ Erro ao atualizar status: {e}")
                last_gif_task = asyncio.run_coroutine_threadsafe(_enviar_status(), loop)

            async def run_mention_task():
                nonlocal msg_ref, last_gif_task

                arquivos, erro_anexo = await ler_anexos_discord(message.attachments)
                if erro_anexo:
                    if msg_ref:
                        await msg_ref.edit(content=erro_anexo)
                    else:
                        await message.reply(erro_anexo)
                    return

                try:
                    resposta, img_gerada = await asyncio.to_thread(
                        self.processar_gemini,
                        conteudo_limpo or "[Enviou um arquivo em anexo sem nenhuma mensagem de texto. Analise o arquivo enviado diretamente (não use a ferramenta ver_tela_atual, pois o arquivo já foi enviado nos bytes do anexo)]",
                        arquivos,
                        "Discord",
                        message,
                        status_callback
                    )
                except Exception as e:
                    resposta = f"Eita, deu um probleminha aqui: {e} 🥺"
                    img_gerada = None

                if last_gif_task:
                    try:
                        await asyncio.wrap_future(last_gif_task)
                    except Exception:
                        pass

                arquivo_discord = None
                if img_gerada and Path(img_gerada).is_file():
                    arquivo_discord = discord.File(img_gerada, filename=Path(img_gerada).name)

                try:
                    pedacos = dividir_mensagem_discord(resposta)
                    if msg_ref:
                        if arquivo_discord:
                            try:
                                await msg_ref.delete()
                            except Exception:
                                pass
                            await message.reply(pedacos[0], file=arquivo_discord)
                        else:
                            await msg_ref.edit(content=pedacos[0])
                    else:
                        await message.reply(pedacos[0], file=arquivo_discord)
                    for pedaco in pedacos[1:]:
                        if pedaco.strip():
                            await message.channel.send(pedaco)
                except discord.errors.Forbidden:
                    print("⚠️ Sem permissão para enviar mensagem ou arquivo na menção.")
                except Exception as e:
                    print(f"⚠️ Erro ao enviar resposta da menção: {e}")

            async with message.channel.typing():
                await run_mention_task()

            return  # Menção tratada

        # ── MODO SOCIALIZAR (Chat Livre em Lotes de 5 em 5) ──
        if message.channel is not None and ayla_state and ayla_state.is_modo_socializar_ativo(message.channel.id):
            msg_info = {
                "author": message.author.display_name,
                "username": message.author.name,
                "author_id": message.author.id,
                "content": message.content.strip() or "[Mídia / Anexo]",
                "timestamp": datetime.now().strftime("%H:%M:%S")
            }
            lote = ayla_state.adicionar_mensagem_socializar(message.channel.id, msg_info)
            if lote:
                asyncio.create_task(self.processar_lote_socializar(message.channel, lote))

    async def processar_lote_socializar(self, channel, lote):
        try:
            mensagens_fmt = []
            for idx, msg in enumerate(lote, 1):
                mensagens_fmt.append(f"[{idx}] {msg['author']} (@{msg['username']}) às {msg['timestamp']}: \"{msg['content']}\"")
            
            bloco_texto = "\n".join(mensagens_fmt)
            prompt_social = (
                "[MODO SOCIALIZAR ATIVO - LEITURA DE CHAT EM LOTE (5/5)]\n"
                "Você está no MODO SOCIALIZAR conversando no chat sem precisar de gatilhos (/ayla ou @ayla).\n"
                "Acabaram de chegar 5 mensagens recentes no chat do servidor:\n\n"
                f"{bloco_texto}\n\n"
                "INSTRUÇÕES PARA A IA:\n"
                "- Analise as 5 mensagens com atenção.\n"
                "- Você pode responder diretamente a uma ou mais pessoas específicas (ex: '@Usuario ...'), fazer um comentário geral sobre a conversa, ou mandar uma mensagem espontânea da Ayla.\n"
                "- Se as mensagens do lote forem spam, irrelevantes ou sem nexo que não merecem resposta, responda APENAS a palavra: [SILENCIO]\n"
            
            )
            
            async with channel.typing():
                resposta, img_gerada = await asyncio.to_thread(
                    self.processar_gemini,
                    prompt_social,
                    [],
                    "ModoSocializar",
                    None,
                    None
                )
            
            if resposta and "[SILENCIO]" not in resposta.upper():
                pedacos = dividir_mensagem_discord(resposta)
                arquivo_discord = None
                if img_gerada and Path(img_gerada).is_file():
                    arquivo_discord = discord.File(img_gerada, filename=Path(img_gerada).name)
                
                await channel.send(pedacos[0], file=arquivo_discord)
                for p in pedacos[1:]:
                    if p.strip():
                        await channel.send(p)
        except Exception as e:
            print(f"⚠️ Erro ao processar lote do Modo Socializar: {e}")

    def enviar_com_fallback(self, payload, modelos_list=None):
        if modelos_list is None:
            modelos_list = self.modelos_disponiveis
        modelos_list = [str(modelo).strip() for modelo in modelos_list if str(modelo).strip()]
        if not modelos_list:
            raise RuntimeError("Nenhum modelo Gemini válido foi configurado.")

        if self.chat_session is None:
            self.modelo_atual = modelos_list[0]
            self.chat_session = self.genai_client.chats.create(
                model=self.modelo_atual,
                config=self.config,
                history=[],
            )

        try:
            idx_modelo = modelos_list.index(self.modelo_atual)
        except ValueError:
            idx_modelo = 0
            self.modelo_atual = modelos_list[0]
            try:
                historico = limitar_historico(self.chat_session.get_history()) if self.chat_session else []
            except Exception:
                historico = []
            self.chat_session = self.genai_client.chats.create(
                model=self.modelo_atual,
                config=self.config,
                history=historico
            )

        while True:
            try:
                t0 = time.time()
                resp = self.chat_session.send_message(payload, config=self.config)
                dur = time.time() - t0
                self.model_status[self.modelo_atual] = {
                    "status": "slow" if dur > 5.0 else "ok",
                    "latency": f"{dur:.1f}s",
                    "latency_ms": int(dur * 1000),
                    "error": ""
                }
                return resp
            except Exception as erro:
                erro_msg = str(erro)
                self.model_status[self.modelo_atual] = {
                    "status": "error",
                    "latency": "N/A",
                    "error": erro_msg[:80]
                }
                erro_msg = str(erro)
                is_quota        = "429" in erro_msg or "RESOURCE_EXHAUSTED" in erro_msg
                is_interno      = "500" in erro_msg or "INTERNAL" in erro_msg
                is_indisponivel = "503" in erro_msg or "UNAVAILABLE" in erro_msg
                is_modality     = "400" in erro_msg and ("modality" in erro_msg.lower() or "not supported" in erro_msg.lower())
                is_403          = "403" in erro_msg or "PERMISSION_DENIED" in erro_msg

                is_suspended = "suspended" in erro_msg.lower() or "consumer" in erro_msg.lower()
                is_403_file = is_403 and not is_suspended

                if is_403_file:
                    print("\n⚠️ Erro 403: Arquivo expirado no Google. Reiniciando sessão...")
                    try:
                        self.chat_session = self.genai_client.chats.create(model=self.modelo_atual, config=self.config)
                        return self.chat_session.send_message(payload, config=self.config)
                    except Exception as sub_err:
                        erro = sub_err
                        erro_msg = str(sub_err)
                        # Re-calcula as flags caso o erro tenha mudado
                        is_quota        = "429" in erro_msg or "RESOURCE_EXHAUSTED" in erro_msg
                        is_interno      = "500" in erro_msg or "INTERNAL" in erro_msg
                        is_indisponivel = "503" in erro_msg or "UNAVAILABLE" in erro_msg
                        is_modality     = "400" in erro_msg and ("modality" in erro_msg.lower() or "not supported" in erro_msg.lower())
                        is_suspended    = "suspended" in erro_msg.lower() or "consumer" in erro_msg.lower()

                if is_interno or is_indisponivel or is_modality:
                    motivo = "Modelo não aceita mídia" if is_modality else ("API sobrecarregada (503)" if is_indisponivel else "Erro 500 interno")
                    print(f"\n⚠️ {motivo}! Pulando de modelo...")
                    idx_modelo += 1
                    if idx_modelo >= len(modelos_list):
                        idx_modelo -= 1
                        raise Exception("Todos os modelos falharam.")

                    self.modelo_atual = modelos_list[idx_modelo]
                    try:
                        historico = limitar_historico(self.chat_session.get_history()) if self.chat_session else []
                    except Exception:
                        historico = []
                    self.chat_session = self.genai_client.chats.create(model=self.modelo_atual, config=self.config, history=historico)
                    continue

                if is_quota or is_suspended:
                    self.idx_api_atual += 1

                    if self.idx_api_atual >= len(self.api_keys):
                        self.idx_api_atual = 0
                        idx_modelo += 1

                        if idx_modelo >= len(modelos_list):
                            raise Exception("Todas as contas e modelos esgotaram a quota ou foram suspensos.")

                        self.modelo_atual = modelos_list[idx_modelo]
                        print(f"\n🔄 Todas as APIs esgotadas/suspensas no modelo anterior. Pulando para {self.modelo_atual}...")
                    else:
                        print(f"\n🔄 Quota/Suspensão de API! Indo para api {self.idx_api_atual + 1}...")

                    self._trocar_cliente_genai(self.api_keys[self.idx_api_atual])

                    try:
                        historico = limitar_historico(self.chat_session.get_history()) if (self.chat_session and hasattr(self.chat_session, "get_history")) else []
                    except Exception:
                        historico = []

                    self.chat_session = self.genai_client.chats.create(
                        model=self.modelo_atual,
                        config=self.config,
                        history=historico
                    )
                else:
                    print(f"\n⚠️ Erro genérico no modelo ({erro}): tentando próximo modelo...")
                    idx_modelo += 1
                    if idx_modelo >= len(modelos_list):
                        idx_modelo -= 1
                        raise erro

                    self.modelo_atual = modelos_list[idx_modelo]
                    try:
                        historico = limitar_historico(self.chat_session.get_history()) if self.chat_session else []
                    except Exception:
                        historico = []
                    self.chat_session = self.genai_client.chats.create(model=self.modelo_atual, config=self.config, history=historico)
                    continue

    def tentar_comando_dm_whitelist(self, user_input: str, origem: str = "Discord") -> str | None:
        # Sobrescrita dinamicamente pelo módulo enviar_mensagem_whitelist_bot.py
        return None

    def processar_gemini(self, user_input: str, arquivos: list[tuple[bytes, str]] = None, origem: str = "Discord", contexto_discord = None, status_callback = None) -> tuple[str, str | None]:
        task_id = str(id(contexto_discord)) if contexto_discord else "local"
        ayla_state.adicionar_tarefa(task_id, f"Processando mensagem em {origem}...")
        try:
            ayla_state.ULTIMA_IMAGEM_GERADA.set(None)
            ayla_state.ULTIMO_ANEXO_IMAGEM.set(None)
            ayla_state.ULTIMO_ANEXO_VIDEO.set(None)
            ayla_state.CONTEXTO_ATIVO.set(contexto_discord)
            ayla_state.ULTIMAS_IMAGENS_MODULO.set([])

            def finalizar_interacao(resposta: str, img=None, status: str = "sucesso"):
                resposta_processada = processar_tags_midia_ayla(resposta)
                return resposta_processada, img

            contexto_privado = origem in ("GUI", "Terminal")
            if arquivos:
                for arq_bytes, tipo in arquivos:
                    if tipo.startswith("image/") and not ayla_state.ULTIMO_ANEXO_IMAGEM.get():
                        ayla_state.ULTIMO_ANEXO_IMAGEM.set((arq_bytes, tipo))
                    elif tipo.startswith("video/") and not ayla_state.ULTIMO_ANEXO_VIDEO.get():
                        pasta_video = BASE_DIR / "Buffer de video"
                        pasta_video.mkdir(parents=True, exist_ok=True)
                        import uuid
                        ext = ".mp4"
                        if "/" in tipo:
                            ext_sug = tipo.split("/")[-1]
                            if ext_sug in ("mp4", "webm", "mkv", "avi", "mov"):
                                ext = f".{ext_sug}"
                        caminho_video = pasta_video / f"video_{uuid.uuid4().hex[:8]}{ext}"
                        caminho_video.write_bytes(arq_bytes)
                        ayla_state.ULTIMO_ANEXO_VIDEO.set(str(caminho_video))

            try:
                # Identifica ID do autor
                autor_id = 0
                ctx = ayla_state.CONTEXTO_ATIVO.get()
                if origem not in ("GUI", "Terminal") and ctx is not None:
                    if hasattr(ctx, "author"):
                        autor_id = getattr(ctx.author, "id", 0)
                    elif hasattr(ctx, "user"):
                        autor_id = getattr(ctx.user, "id", 0)

                self.ayla_settings = carregar_configuracoes_ayla()
                
                # Ajusta dinamicamente a instrução do prompt em relação ao limite de fala/áudio
                sys_prompt_atual = self.prompt_com_memoria
                if self.ayla_settings.get("tts_engine") == "fish_audio":
                    sys_prompt_atual = sys_prompt_atual.replace(
                        "- Como a sua fala de áudio é limitada a apenas 50 caracteres (letras), escreva textos bem curtinhos e diretos quando for usar a ferramenta de áudio!",
                        "- Como você está usando o Fish Audio, não há limite de caracteres para sua fala. Escreva livremente!"
                    )
                self.config.system_instruction = sys_prompt_atual
                self.aparar_historico_chat()

                resultado_dm = self.tentar_comando_dm_whitelist(user_input, origem)
                if resultado_dm is not None:
                    # Clear globals before returning
                    ayla_state.ULTIMA_IMAGEM_GERADA.set(None)
                    ayla_state.CONTEXTO_ATIVO.set(None)
                    ayla_state.ULTIMO_ANEXO_IMAGEM.set(None)
                    ayla_state.ULTIMO_ANEXO_VIDEO.set(None)
                    ayla_state.ULTIMAS_IMAGENS_MODULO.set([])
                    ayla_state.remover_tarefa(task_id)
                    return finalizar_interacao(resultado_dm, None, "sucesso")


                if origem in ("GUI", "Terminal"):
                    usuario_nome = "User"
                    local_info = origem
                else:
                    usuario_nome = "Usuário"
                    local_info = "Discord"
                    ctx = ayla_state.CONTEXTO_ATIVO.get()
                    if ctx:
                        if hasattr(ctx, "author"):
                            usuario_nome = str(getattr(ctx.author, "display_name", ctx.author))
                        elif hasattr(ctx, "user"):
                            usuario_nome = str(getattr(ctx.user, "display_name", ctx.user))

                        guild_obj = getattr(ctx, "guild", None)
                        channel_obj = getattr(ctx, "channel", None)
                        canal_nome = getattr(channel_obj, "name", "") if channel_obj else ""

                        if guild_obj:
                            srv_nome = guild_obj.name
                            c_info = f" | Canal #{canal_nome}" if canal_nome else ""
                            if hasattr(ctx, "command") or type(ctx).__name__ == "Interaction":
                                local_info = f"Comando Slash no Servidor '{srv_nome}'{c_info}"
                            else:
                                local_info = f"Servidor '{srv_nome}'{c_info}"
                        else:
                            local_info = "Mensagem Direta (DM)"

                agora = datetime.now().strftime("%d/%m/%Y %H:%M:%S")
                prefixo = f"[Data/Hora: {agora}] [local: {local_info}] [nome: {usuario_nome}] [id: {autor_id}]"
                texto_formatado = f"{prefixo}: {user_input}" if user_input else f"{prefixo} enviou um arquivo"

                modelos_req = list(getattr(self, "modelos_padrao", []))
                if not modelos_req:
                    modelos_req = ["gemini-3.1-flash-lite", "gemini-2.5-flash-lite"]

                target_model = modelos_req[0]
                if getattr(self, "modelo_atual", None) != target_model:
                    print(f"🔄 Alternando modelo da sessão de {getattr(self, 'modelo_atual', None)} para {target_model}...")
                    try:
                        historico = limitar_historico(self.chat_session.get_history()) if self.chat_session else []
                    except Exception:
                        historico = []
                    self.chat_session = self.genai_client.chats.create(
                        model=target_model,
                        config=self.config,
                        history=historico
                    )
                    self.modelo_atual = target_model

                payload = []
                textos_documentos = []

                if arquivos:
                    for arq_bytes, arq_mime in arquivos:
                        if arq_mime.startswith("text/document:"):
                            # Documento de texto: injeta conteúdo como texto no contexto
                            nome_doc = arq_mime.split("text/document:", 1)[1]
                            conteudo_doc = arq_bytes.decode("utf-8", errors="replace")
                            textos_documentos.append(
                                f"📄 **Conteúdo do arquivo '{nome_doc}':**\n```\n{conteudo_doc}\n```"
                            )
                        else:
                            mime_limpo = arq_mime.split(";")[0]
                            payload.append(genai_types.Part.from_bytes(data=arq_bytes, mime_type=mime_limpo))

                # Monta o texto final combinando documentos + mensagem do usuário
                texto_final = texto_formatado
                if textos_documentos:
                    cabecalho_docs = "\n\n".join(textos_documentos)
                    texto_final = f"{cabecalho_docs}\n\n{texto_formatado}"

                if payload:
                    payload.append(genai_types.Part(text=texto_final))
                else:
                    payload = texto_final

                response = self.enviar_com_fallback(payload, modelos_list=modelos_req)

                textos_acumulados = []

                loop_seguranca = 0
                while loop_seguranca < 20:
                    loop_seguranca += 1

                    if not response.candidates or not response.candidates[0].content or not response.candidates[0].content.parts: break

                    # Separa partes de texto e partes de function_call
                    fn_calls = []
                    for p in response.candidates[0].content.parts:
                        if p.function_call and p.function_call.name:
                            fn_calls.append(p)
                        elif hasattr(p, "text") and p.text and p.text.strip():
                            textos_acumulados.append(p.text.strip())

                    if not fn_calls: break

                    # Verifica se quem chamou NÃO é a dona (modo público)
                    _is_owner = (autor_id == DISCORD_OWNER_ID) or origem in ("GUI", "Terminal")

                    fn_responses = []
                    for part in fn_calls:
                        fc = part.function_call
                        nome = fc.name
                        args = dict(fc.args) if fc.args else {}
                        print(f"💻 [{origem}] Executando: {nome}({args})")

                        if not _is_owner and nome != "bloquear_usuario":
                            print(f"🚫 Ferramenta '{nome}' bloqueada — usuário não é a dona.")
                            resultado = ("⚠️ Essa função foi chamada por uma pessoa que não é sua dona. "
                                         "Se foi solicitado pela pessoa, dê uma bronca fofa dizendo que "
                                         "só sua dona pode usar suas ferramentas!")
                        else:
                            if status_callback:
                                try:
                                    status_callback(f"Executando: **{nome}** ({args})\n{GIF_GERAR_IMAGEM}")
                                except Exception as e_sc:
                                    print(f"⚠️ Erro ao chamar status_callback: {e_sc}")

                            resultado = executar_ferramenta(nome, args)
                        
                        res_str = str(resultado)
                        if len(res_str) > 10000:
                            res_exibir = res_str[:10000] + "\n... [Resultado truncado no log do terminal por tamanho excessivo]"
                        else:
                            res_exibir = res_str
                        
                        if "\n" in res_exibir:
                            print(f"   ➡️ Resultado:\n{res_exibir}")
                        else:
                            print(f"   ➡️ Resultado: {res_exibir}")

                        resultado_para_modelo = str(resultado)
                        if len(resultado_para_modelo) > MAX_RESULTADO_FERRAMENTA_CARACTERES:
                            resultado_para_modelo = (
                                resultado_para_modelo[:MAX_RESULTADO_FERRAMENTA_CARACTERES]
                                + "\n[Resultado truncado por segurança.]"
                            )

                        fn_responses.append(
                            genai_types.Part(
                                function_response=genai_types.FunctionResponse(
                                    name=nome,
                                    id=getattr(fc, "id", None),
                                    response={"result": resultado_para_modelo}
                                )
                            )
                        )

                    ultimas_imagens = ayla_state.ULTIMAS_IMAGENS_MODULO.get()
                    if ultimas_imagens:
                        total_imagens = 0
                        imagens_validas = []
                        for img_bytes, mime in ultimas_imagens[:MAX_ANEXOS_POR_REQUISICAO]:
                            if len(img_bytes) > MAX_BYTES_POR_ANEXO:
                                continue
                            if total_imagens + len(img_bytes) > MAX_BYTES_ANEXOS_TOTAL:
                                break
                            total_imagens += len(img_bytes)
                            imagens_validas.append((img_bytes, mime))
                        for img_bytes, mime in imagens_validas:
                            mime_limpo = mime.split(";")[0]
                            fn_responses.append(genai_types.Part.from_bytes(data=img_bytes, mime_type=mime_limpo))
                        if imagens_validas:
                            ayla_state.ULTIMO_ANEXO_IMAGEM.set(imagens_validas[-1])
                        print(f"   ➡️ Injetado(s) {len(imagens_validas)} anexo(s) de imagem na resposta do Gemini.")
                        ayla_state.ULTIMAS_IMAGENS_MODULO.set([])

                    response = self.enviar_com_fallback(fn_responses, modelos_list=modelos_req)



                # Extrai texto da resposta final (após todas as ferramentas)
                texto_resposta_final = ""
                try:
                    if response and response.candidates and response.candidates[0].content and response.candidates[0].content.parts:
                        partes_texto = [p.text.strip() for p in response.candidates[0].content.parts if hasattr(p, "text") and p.text and p.text.strip()]
                        texto_resposta_final = "\n".join(partes_texto)
                except Exception:
                    try:
                        texto_resposta_final = response.text.strip() if response and hasattr(response, "text") and response.text else ""
                    except Exception:
                        texto_resposta_final = ""

                chamadas_pendentes = False
                try:
                    chamadas_pendentes = any(
                        getattr(part, "function_call", None)
                        for part in response.candidates[0].content.parts
                    )
                except Exception:
                    chamadas_pendentes = False

                status_final = "sucesso"
                # Combina: prioriza o texto da resposta final; se vazio, usa textos acumulados durante as chamadas de ferramenta
                if chamadas_pendentes and loop_seguranca >= 20:
                    texto_final = (
                        "Não consegui concluir a tarefa porque ela ultrapassou "
                        "o limite seguro de 20 rodadas de ferramentas."
                    )
                    status_final = "erro"
                elif texto_resposta_final:
                    texto_final = texto_resposta_final
                elif textos_acumulados:
                    texto_final = "\n".join(textos_acumulados)
                else:
                    texto_final = "Não recebi uma resposta válida do modelo. Tente novamente."
                    status_final = "erro"

                img_gerada = ayla_state.ULTIMA_IMAGEM_GERADA.get()
                
                # Clear globals on success
                ayla_state.ULTIMA_IMAGEM_GERADA.set(None)
                ayla_state.CONTEXTO_ATIVO.set(None)
                ayla_state.ULTIMO_ANEXO_IMAGEM.set(None)
                ayla_state.ULTIMO_ANEXO_VIDEO.set(None)
                ayla_state.ULTIMAS_IMAGENS_MODULO.set([])
                ayla_state.remover_tarefa(task_id)

                return finalizar_interacao(texto_final, img_gerada, status_final)

            except Exception as e:
                import traceback
                traceback.print_exc()
                
                # Clear globals on error
                ayla_state.ULTIMA_IMAGEM_GERADA.set(None)
                ayla_state.CONTEXTO_ATIVO.set(None)
                ayla_state.ULTIMO_ANEXO_IMAGEM.set(None)
                ayla_state.ULTIMO_ANEXO_VIDEO.set(None)
                ayla_state.ULTIMAS_IMAGENS_MODULO.set([])
                ayla_state.remover_tarefa(task_id)
                
                erro_final = (
                    f"Vish, capotei feio aqui: {e}"
                    if contexto_privado
                    else "Vish, ocorreu um erro interno enquanto eu processava a solicitação."
                )
                return finalizar_interacao(erro_final, None, "erro")

        finally:
            vid_temp = ayla_state.ULTIMO_ANEXO_VIDEO.get()
            if vid_temp and os.path.isfile(vid_temp):
                try:
                    os.remove(vid_temp)
                except Exception as e_rm:
                    print(f"⚠️ Erro ao apagar vídeo temporário: {e_rm}")
            ayla_state.remover_tarefa(task_id)

bot = AylaBot()
carregar_modulos()
bot.inicializar_configuracao_gemini()
_sincronizar_contexto_modulos()


# ══════════════════════════════════════════════════════════
#  HELPER DE ENVIO
# ══════════════════════════════════════════════════════════

async def _enviar_resposta(interaction: discord.Interaction, resposta: str, arquivo_discord=None, edit_original: bool = False):
    pedacos = dividir_mensagem_discord(resposta)
    if edit_original:
        if arquivo_discord:
            try:
                await interaction.delete_original_response()
            except Exception as e:
                print(f"⚠️ Erro ao deletar original: {e}")
        else:
            try:
                await interaction.edit_original_response(content=pedacos[0])
                for pedaco in pedacos[1:]:
                    if pedaco.strip():
                        await interaction.channel.send(pedaco)
                return
            except Exception as e:
                print(f"⚠️ Erro ao editar original, enviando followup: {e}")

    kwargs = {"file": arquivo_discord} if arquivo_discord else {}
    try:
        await interaction.followup.send(pedacos[0], **kwargs)
        for pedaco in pedacos[1:]:
            if pedaco.strip():
                await interaction.channel.send(pedaco)
    except discord.errors.Forbidden:
        print("⚠️ Sem permissão (403).")
    except Exception as e:
        print(f"Erro ao mandar resposta: {e}")


# ══════════════════════════════════════════════════════════
#  COMANDOS DISCORD
# ══════════════════════════════════════════════════════════

# ── Helper de Processamento e Resposta ──────────────────
async def processar_e_enviar_resposta(interaction: discord.Interaction, mensagem: str, anexos_list: list):
    await interaction.response.defer(ephemeral=False)

    msg_ref_gif = False
    last_gif_task = None
    loop = asyncio.get_running_loop()

    def status_callback(status_text: str):
        nonlocal msg_ref_gif, last_gif_task
        async def _enviar_status():
            nonlocal msg_ref_gif
            try:
                await interaction.edit_original_response(content=status_text)
                msg_ref_gif = True
            except discord.errors.NotFound:
                pass
            except Exception as e:
                print(f"⚠️ Erro ao enviar status em slash command: {e}")

        last_gif_task = asyncio.run_coroutine_threadsafe(_enviar_status(), loop)

    arquivos, erro_anexo = await ler_anexos_discord(anexos_list)
    if erro_anexo:
        await interaction.followup.send(erro_anexo)
        return

    texto_envio = mensagem if mensagem else "[Enviou um arquivo em anexo sem nenhuma mensagem de texto. Analise o arquivo enviado diretamente (não use a ferramenta ver_tela_atual, pois o arquivo já foi enviado nos bytes do anexo)]"

    try:
        resposta, img_gerada = await asyncio.to_thread(
            bot.processar_gemini,
            texto_envio,
            arquivos,
            "Discord",
            interaction,
            status_callback
        )
    except Exception as e:
        await interaction.followup.send(f"⚠️ Erro inesperado: {e}")
        return

    if last_gif_task:
        try:
            await asyncio.wrap_future(last_gif_task)
        except Exception:
            pass

    arquivo_discord = None
    if img_gerada and Path(img_gerada).is_file():
        arquivo_discord = discord.File(img_gerada, filename=Path(img_gerada).name)

    await _enviar_resposta(interaction, resposta, arquivo_discord, edit_original=True)


# ── /ayla ────────────────────────────────────────────────
@bot.tree.command(name="ayla", description="Converse com a Ayla — controle total do PC")
@app_commands.describe(
    mensagem="O que quer que eu faça?",
    anexo="Anexo (imagem, áudio, vídeo, PDF, DOCX, TXT, CSV, XLSX...)"
)
@app_commands.allowed_installs(guilds=True, users=True)
@app_commands.allowed_contexts(guilds=True, dms=True, private_channels=True)
async def cmd_ayla(
    interaction: discord.Interaction,
    mensagem: str = "",
    anexo: discord.Attachment = None
):
    if DISCORD_OWNER_ID and interaction.user.id != DISCORD_OWNER_ID:
        _settings = carregar_configuracoes_ayla()
        if not _settings.get("public_mode", False):
            await interaction.response.send_message("⚠️ Apenas a minha dona oficial pode falar comigo!", ephemeral=True)
            return
        if _usuario_bloqueado(interaction.user.id):
            await interaction.response.send_message("⚠️ Você está bloqueado e não pode falar comigo!", ephemeral=True)
            return

    anexos_list = [anexo] if anexo is not None else []
    if not mensagem.strip() and not anexos_list:
        await interaction.response.send_message("⚠️ Fala alguma coisa comigo primeiro ou mande um anexo! https://media.tenor.com/DyqH3PQFYpsAAAAC/burn.gif ", ephemeral=True)
        return

    await processar_e_enviar_resposta(interaction, mensagem, anexos_list)



# ── Bot Discord em thread separada ──
def run_discord_bot_thread():
    import threading
    def _run():
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            async def _start():
                async with bot:
                    await bot.start(DISCORD_TOKEN)
            loop.run_until_complete(_start())
        except Exception as e:
            ayla_state.BOT_LAUNCH_ERROR.set(str(e))
            print(f"\n⚠️ Erro crítico ao iniciar bot Discord: {e}\n")

    discord_thread = threading.Thread(target=_run, daemon=True)
    discord_thread.start()
    print ("Bot Discord iniciando em background...")


if __name__ == "__main__":
    import signal
    signal.signal(signal.SIGINT, lambda *_: (print("\n🩵 Ctrl+C — Ayla encerrando..."), os._exit(0)))

    from colorama import Fore, Back, Style
    print("\n" + "═" * 60)
    print(r"🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵")
    print(Fore.CYAN + r" .----------------. .----------------. .----------------. .----------------.")
    print(r"| .--------------. | .--------------. | .--------------. | .--------------. |")
    print(r"| |      __      | | |  ____  ____  | | |   _____      | | |      __      | |")
    print(r"| |     /  \     | | | |_  _||_  _| | | |  |_   _|     | | |     /  \     | |")
    print(r"| |    / /\ \    | | |   \ \  / /   | | |    | |       | | |    / /\ \    | |")
    print(r"| |   / ____ \   | | |    \ \/ /    | | |    | |   _   | | |   / ____ \   | |")
    print(r"| | _/ /    \ \_ | | |    _|  |_    | | |   _| |__/ |  | | | _/ /    \ \_ | |")
    print(r"| ||____|  |____|| | |   |______|   | | |  |________|  | | ||____|  |____|| |")
    print(r"| |              | | |              | | |              | | |              | |")
    print(r"| '--------------' | '--------------' | '--------------' | '--------------' |")
    print(r"'----------------' '----------------' '----------------' '------------------'")
    print(Style.RESET_ALL)
    print(r"🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵🩵")



    print("═" * 60 + "\n")

    abrir_gui_env = os.getenv("ABRIR_GUI", "true").lower() == "true"
    if abrir_gui_env:
        run_discord_bot_thread()
        try:
            from ayla_gui import abrir_gui
            print("🩵 Abrindo interface gráfica...")
            abrir_gui(bot)
            # Mantém a thread principal viva após fechar a GUI
            import time
            while True:
                time.sleep(1)
        except KeyboardInterrupt:
            print("\n🩵 Ctrl+C — Ayla encerrando...")
            os._exit(0)
        except Exception as e:
            print(f"⚠️ Erro ao abrir GUI: {e}")
            print("🩵 Iniciando bot no thread principal...")
            try:
                async def main():
                    async with bot:
                        await bot.start(DISCORD_TOKEN)
                asyncio.run(main())
            except KeyboardInterrupt:
                print("\n🩵 Ctrl+C — Ayla encerrando...")
                os._exit(0)
    else:
        print("🩵 Iniciando bot no thread principal...")
        try:
            async def main():
                async with bot:
                    await bot.start(DISCORD_TOKEN)
            asyncio.run(main())
        except KeyboardInterrupt:
            print("\n🩵 Ctrl+C — Ayla encerrando...")
            os._exit(0)

